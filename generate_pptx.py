#!/usr/bin/env python3
"""
generate_pptx.py  —  Step 6: Generate the PowerPoint deck.

Reads the issues markdown, selections.json, and monitor_template.pptx,
then produces one slide per issue (in severity order) with text and
the selected/cropped screenshot filled into the named shapes.

Named shapes required in the template slide:
    issue_id        severity_badge      title
    observed        expected            notes      (notes may be absent)
    metadata        screenshot

Usage:
    python generate_pptx.py
"""

import copy
import json
import os
import re
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

CONFIG_FILE = "config.json"

SEVERITY_COLORS = {
    "S0": RGBColor(0xC0, 0x00, 0x00),  # red
    "S1": RGBColor(0xC5, 0x50, 0x00),  # red-orange
    "S2": RGBColor(0xE0, 0x80, 0x00),  # orange
    "S3": RGBColor(0x99, 0x7A, 0x00),  # burnt yellow
}


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

def parse_issues(md_path):
    with open(md_path, encoding="utf-8") as f:
        content = f.read()

    sections = re.split(r"(?m)(?=^# VID-)", content)
    issues = []

    for section in sections:
        section = section.strip()
        if not section:
            continue

        lines = section.split("\n")
        header = re.match(r"# (VID-\d+)\s*[-\u2013]\s*(.+)", lines[0])
        if not header:
            continue

        issue = {
            "id":             header.group(1),
            "title":          header.group(2).strip(),
            "categories":     "",
            "severity":       "",
            "confidence":     "",
            "timestamps":     "",
            "affected_roles": "",
            "affected_area":  "",
            "observed":       "",
            "expected":       "",
            "notes":          "",
        }

        for line in lines[1:20]:
            m = re.match(r"-\s+\*\*(.+?):\*\*\s*(.*)", line)
            if m:
                key = m.group(1).strip().lower()
                val = m.group(2).strip()
                mapping = {
                    "categories":     "categories",
                    "severity":       "severity",
                    "confidence":     "confidence",
                    "timestamps":     "timestamps",
                    "affected roles": "affected_roles",
                    "affected area":  "affected_area",
                }
                if key in mapping:
                    issue[mapping[key]] = val

        for field, heading in [
            ("observed", "Observed behavior"),
            ("expected", "Expected behavior"),
            ("notes",    "Notes"),
        ]:
            pat = rf"## {re.escape(heading)}\n(.*?)(?=\n## |\Z)"
            m = re.search(pat, section, re.DOTALL | re.IGNORECASE)
            if m:
                issue[field] = m.group(1).strip()

        issues.append(issue)

    return issues


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def duplicate_slide(prs, template_idx=0):
    """Duplicate the slide at template_idx and append it to the presentation."""
    template = prs.slides[template_idx]
    slide_layout = template.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    sp_tree = new_slide.shapes._spTree
    # Clear every child the layout added
    for child in list(sp_tree):
        sp_tree.remove(child)
    # Copy all children from the template's shape tree
    for child in template.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))

    return new_slide


def get_shape(slide, name):
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def _write_run_text(para, text):
    """Set text on a paragraph's first run, preserving its font properties.
    Removes any extra runs.

    When no runs exist (empty template shape), the font settings live in the
    paragraph's a:endParaRPr element.  We clone that into the new run's a:rPr
    so that font size, color, bold, etc. set in the template are honoured.
    """
    if para.runs:
        para.runs[0].text = text
        for run in para.runs[1:]:
            run._r.getparent().remove(run._r)
        return

    # No existing run — create one via python-pptx (inserts before endParaRPr)
    run = para.add_run()
    run.text = text

    # Copy template formatting from endParaRPr → new run's rPr
    end_rpr = para._p.find(qn('a:endParaRPr'))
    if end_rpr is not None:
        rpr = copy.deepcopy(end_rpr)
        rpr.tag = qn('a:rPr')
        run._r.insert(0, rpr)  # a:rPr must be the first child of a:r


def set_font_color(slide, shape_name, rgb_color):
    """Override the font color on every run in a shape."""
    shape = get_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb_color


def _remove_extra_paragraphs(tf):
    """Remove all paragraphs after the first from a text frame."""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)


def set_text(slide, shape_name, text):
    shape = get_shape(slide, shape_name)
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    _remove_extra_paragraphs(tf)
    _write_run_text(tf.paragraphs[0], text)


def set_bullets(slide, shape_name, markdown_text):
    """
    Fill a text frame with bullet paragraphs parsed from markdown.
    Top-level bullets (- text) → level 0
    Indented sub-bullets (    - text) → level 1
    Preserves the template's run formatting (font size, bold, etc.).
    """
    shape = get_shape(slide, shape_name)
    if not shape or not shape.has_text_frame or not markdown_text:
        return

    lines = []
    for line in markdown_text.split("\n"):
        stripped = line.lstrip()
        if not stripped:
            continue
        # Skip transcript quotes (sub-bullets whose text starts with a quote mark)
        bullet_text = stripped[2:] if stripped.startswith("- ") else stripped
        if bullet_text.startswith('"'):
            continue
        indent = len(line) - len(stripped)
        lines.append((bullet_text, 1 if indent >= 4 else 0))

    if not lines:
        return

    tf = shape.text_frame
    # Save a copy of the template paragraph XML to clone for each new paragraph
    para_template = copy.deepcopy(tf.paragraphs[0]._p)
    _remove_extra_paragraphs(tf)

    # Fill the first paragraph in-place
    _write_run_text(tf.paragraphs[0], lines[0][0])
    tf.paragraphs[0].level = lines[0][1]

    # Append cloned paragraphs for the rest, preserving template formatting
    for text, level in lines[1:]:
        new_p = copy.deepcopy(para_template)
        tf._txBody.append(new_p)
        _write_run_text(tf.paragraphs[-1], text)
        tf.paragraphs[-1].level = level


def insert_screenshot(slide, shape_name, image_path, crop=None):
    """
    Replace the named placeholder shape with the screenshot image.
    Applies crop (pixel coords relative to original) before inserting.
    """
    shape = get_shape(slide, shape_name)
    if not shape:
        return

    left   = shape.left
    top    = shape.top
    width  = shape.width
    height = shape.height

    # Remove placeholder
    shape._element.getparent().remove(shape._element)

    tmp_path = None
    try:
        if crop:
            img     = Image.open(image_path)
            cropped = img.crop((crop["left"], crop["top"], crop["right"], crop["bottom"]))
            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
                cropped.save(tmp.name, "JPEG", quality=95)
                tmp_path = tmp.name
            image_path = tmp_path

        # Fit image inside shape bounds while preserving aspect ratio (letterbox/pillarbox)
        img_w, img_h = Image.open(image_path).size
        scale = min(width / img_w, height / img_h)
        fit_w = int(img_w * scale)
        fit_h = int(img_h * scale)
        fit_left = left + (width  - fit_w) // 2
        fit_top  = top  + (height - fit_h) // 2
        slide.shapes.add_picture(image_path, fit_left, fit_top, fit_w, fit_h)
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def render_slide(prs, issue, sel):
    """Duplicate the template slide and populate it from issue + sel overrides."""
    slide = duplicate_slide(prs, template_idx=0)
    vid   = issue["id"]

    # Apply text overrides from sel onto a working copy of the issue fields
    TEXT_FIELDS = ("title", "severity", "observed", "expected", "notes")
    data = dict(issue)
    if sel and isinstance(sel, dict):
        for f in TEXT_FIELDS:
            if f in sel:
                data[f] = sel[f]

    # ── Text fields ──
    set_text(slide, "issue_id",       vid)
    sev_key = data["severity"][:2] if data["severity"] else ""
    set_text(slide, "severity_badge", sev_key)
    if sev_key in SEVERITY_COLORS:
        set_font_color(slide, "severity_badge", SEVERITY_COLORS[sev_key])
    set_text(slide, "title", data["title"])

    set_bullets(slide, "observed", data["observed"])
    set_bullets(slide, "expected", data["expected"])
    if data["notes"]:
        set_bullets(slide, "notes", data["notes"])

    metadata_override = sel.get("metadata") if sel and isinstance(sel, dict) else None
    if metadata_override:
        metadata = metadata_override
    else:
        metadata = (
            f"Roles: {issue['affected_roles']} | "
            f"Area: {issue['affected_area']} | "
            f"{issue['timestamps']}"
        ).strip(" |")
    set_text(slide, "metadata", metadata)

    # ── Screenshot ──
    if sel and isinstance(sel, dict) and sel.get("path"):
        img_path = sel["path"]
        crop     = sel.get("crop")
        if os.path.exists(img_path):
            insert_screenshot(slide, "screenshot", img_path, crop)
        else:
            print(f"    [WARN] Image not found: {img_path}")


def main():
    if not os.path.exists(CONFIG_FILE):
        sys.exit(f"Config not found: {CONFIG_FILE}")

    with open(CONFIG_FILE) as f:
        config = json.load(f)

    issues_path    = config["issues_path"]
    selections_path = config["selections_path"]
    template_path  = config["template_path"]
    output_path    = config["output_pptx"]

    for path in (issues_path, template_path):
        if not os.path.exists(path):
            sys.exit(f"File not found: {path}")

    if not os.path.exists(selections_path):
        print(f"[WARN] {selections_path} not found — all screenshots will be blank.")
        selections = {}
    else:
        with open(selections_path, encoding="utf-8") as f:
            selections = json.load(f)

    issues = parse_issues(issues_path)
    print(f"Parsed {len(issues)} issues from {issues_path}")

    prs = Presentation(template_path)
    if len(prs.slides) == 0:
        sys.exit("Template has no slides.")

    print("Generating slides...")

    slide_count = 0
    for issue in issues:
        vid = issue["id"]
        sel = selections.get(vid) or {}

        if sel.get("merged_into"):
            print(f"  [SKIP] {vid} — merged into {sel['merged_into']}")
            continue

        print(f"  [{slide_count+1:02d}] {vid}  {issue['title'][:60]}")
        render_slide(prs, issue, sel)
        slide_count += 1

        slide_b = sel.get("slide_b")
        if slide_b and isinstance(slide_b, dict):
            merged_b = {k: v for k, v in sel.items() if k != "slide_b"}
            merged_b.update(slide_b)
            print(f"  [{slide_count+1:02d}] {vid}b (split slide)")
            render_slide(prs, issue, merged_b)
            slide_count += 1

    # Remove the original template slide (index 0)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    prs.save(output_path)
    print(f"\nSaved: {output_path}  ({slide_count} slides)")


if __name__ == "__main__":
    main()
